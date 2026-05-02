from __future__ import annotations
import io
import json
import re
import zipfile
from services import ai_service as ai, db_service as db, search_service
from services.auth_service import get_supabase

IMPORT_SYSTEM = """You are ProjectVault's AI project importer. Analyze the provided content and extract project information. Return ONLY a valid JSON object with exactly these keys:
{
  "project_name": "Inferred or explicit project name (string)",
  "description": "One clear paragraph summarizing what this project is about, its goals, and current state",
  "tags": ["tag1", "tag2", "tag3"],
  "milestones": [{"title": "Milestone name", "date": "YYYY-MM-DD or month/year or unknown", "status": "completed|pending|unknown"}],
  "import_narrative": "2-3 sentences about what this project is and what stage it appears to be at",
  "detected_tools": [],
  "health_score": 70,
  "health_explanation": "One sentence explanation of the score based on how complete and clear the documentation is"
}

For detected_tools, ONLY include values from this exact list: claude, chatgpt, notion, figma, jira, slack, googledocs, loom, linear, github.
Health score should reflect documentation completeness: sparse notes = 40-55, decent brief = 55-70, detailed PRD = 70-85, comprehensive = 85-100.
Return ONLY valid JSON, no markdown wrapper."""


# ── Extension sets ────────────────────────────────────────────────────────────

_CODE_EXTS = frozenset({"py", "js", "ts", "jsx", "tsx", "sh", "rb", "go", "rs", "java", "cpp", "c", "swift"})
_TEXT_EXTS = frozenset({"txt", "md", "json", "yaml", "yml", "html", "css", "sql", "toml", "ini", "cfg", "env", "ipynb"})
_BINARY_EXTS = frozenset({
    "png", "jpg", "jpeg", "gif", "svg", "ico", "webp", "bmp",
    "mp4", "mp3", "wav", "mov", "avi", "mkv",
    "zip", "tar", "gz", "7z", "rar",
    "exe", "dll", "so", "dylib", "bin",
    "ttf", "woff", "woff2", "eot",
    "pyc", "pyo", "class",
})
_PRIORITY_FILES = frozenset({"readme.md", "readme.txt", "readme.rst", "requirements.txt", "package.json", "pyproject.toml"})


# ── Text extraction ────────────────────────────────────────────────────────────

def extract_text(file_bytes: bytes, filename: str) -> str:
    ext = filename.lower().rsplit(".", 1)[-1]

    if ext in _TEXT_EXTS | _CODE_EXTS | {"txt", "md"}:
        return file_bytes.decode("utf-8", errors="replace")

    if ext == "pdf":
        try:
            import pdfplumber
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                return "\n".join(page.extract_text() or "" for page in pdf.pages)
        except Exception:
            return ""

    if ext == "docx":
        try:
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception:
            return ""

    if ext == "pptx":
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text)
            return "\n".join(parts)
        except Exception:
            return ""

    if ext in ("xlsx", "csv"):
        try:
            import pandas as pd
            df = pd.read_csv(io.BytesIO(file_bytes)) if ext == "csv" else pd.read_excel(io.BytesIO(file_bytes))
            return df.to_string(index=False)
        except Exception:
            return ""

    return ""


def fetch_url_content(url: str) -> str:
    try:
        import requests
        from bs4 import BeautifulSoup
        resp = requests.get(url, headers={"User-Agent": "Mozilla/5.0 ProjectVault/1.0"}, timeout=10)
        resp.raise_for_status()
        soup = BeautifulSoup(resp.text, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)[:8000]
    except Exception as e:
        return f"Could not fetch URL: {e}"


# ── ZIP extraction ────────────────────────────────────────────────────────────

def extract_text_from_zip(zip_bytes: bytes) -> dict:
    """
    Extract text from all supported files inside a ZIP archive.

    Returns:
        combined_text  – concatenated text ready for AI analysis
        file_tree      – list of (display_name, size_kb, status)
        summary        – human-readable "Found 12 files: 3 PDFs, ..."
        is_code_project – True if 2+ code files detected
        error          – "corrupted" | "password_protected" | None
        truncated      – True if text was cut at 100 KB
    """
    result: dict = {
        "combined_text": "", "file_tree": [], "summary": "",
        "is_code_project": False, "error": None, "truncated": False,
    }

    try:
        zf_obj = zipfile.ZipFile(io.BytesIO(zip_bytes))
    except zipfile.BadZipFile:
        result["error"] = "corrupted"
        return result

    with zf_obj as zf:
        # Detect password protection by trying to read the first file
        entries = [e for e in zf.infolist() if not e.filename.endswith("/")]
        if entries:
            try:
                zf.read(entries[0].filename)
            except RuntimeError:
                result["error"] = "password_protected"
                return result

        priority_parts: list[str] = []
        regular_parts: list[str] = []
        file_counts: dict[str, int] = {}
        code_file_count = 0

        # Sort: priority files first
        def _sort_key(e: zipfile.ZipInfo) -> int:
            return 0 if e.filename.lower().rsplit("/", 1)[-1] in _PRIORITY_FILES else 1

        for entry in sorted(entries, key=_sort_key):
            name = entry.filename
            display = name.rsplit("/", 1)[-1] if "/" in name else name
            size_kb = entry.file_size / 1024

            # Skip hidden / system files
            if display.startswith(".") or display.startswith("__"):
                continue

            ext = display.lower().rsplit(".", 1)[-1] if "." in display else ""

            if ext in _BINARY_EXTS:
                result["file_tree"].append((display, size_kb, "skipped"))
                continue

            try:
                raw = zf.read(name)
            except Exception:
                result["file_tree"].append((display, size_kb, "error"))
                continue

            # Extract text by type
            if ext in _CODE_EXTS | _TEXT_EXTS | {"txt", "md"}:
                text = raw.decode("utf-8", errors="replace")
                category = "code" if ext in _CODE_EXTS else ext
                if ext in _CODE_EXTS:
                    code_file_count += 1
            elif ext == "pdf":
                text = extract_text(raw, display)
                category = "pdf"
            elif ext == "docx":
                text = extract_text(raw, display)
                category = "docx"
            elif ext == "pptx":
                text = extract_text(raw, display)
                category = "pptx"
            elif ext in ("csv", "xlsx"):
                text = extract_text(raw, display)
                category = "spreadsheet"
            else:
                result["file_tree"].append((display, size_kb, "skipped"))
                continue

            if not text.strip():
                result["file_tree"].append((display, size_kb, "empty"))
                continue

            file_counts[category] = file_counts.get(category, 0) + 1
            result["file_tree"].append((display, size_kb, "ok"))

            chunk = f"=== {display} ===\n{text[:4000]}"
            is_priority = display.lower() in _PRIORITY_FILES
            (priority_parts if is_priority else regular_parts).append(chunk)

        result["is_code_project"] = code_file_count >= 2

        combined = "\n\n".join(priority_parts + regular_parts)
        MAX_CHARS = 100_000
        if len(combined) > MAX_CHARS:
            combined = combined[:MAX_CHARS]
            result["truncated"] = True

        result["combined_text"] = combined

        # Human-readable summary
        label_map = {
            "pdf": "PDF", "docx": "Word doc", "pptx": "slide deck",
            "spreadsheet": "spreadsheet", "code": "code file",
            "md": "Markdown", "json": "JSON", "yaml": "YAML",
        }
        parts = [
            f"{n} {label_map.get(cat, cat)}{'s' if n > 1 else ''}"
            for cat, n in file_counts.items()
        ]
        skipped = sum(1 for _, _, s in result["file_tree"] if s == "skipped")
        result["summary"] = (
            f"Found {sum(file_counts.values())} readable file(s): {', '.join(parts)}"
            + (f" · {skipped} skipped (images/binaries)" if skipped else "")
            + (" · large project, analyzing first 100 KB" if result["truncated"] else "")
        )

    return result


# ── AI analysis ────────────────────────────────────────────────────────────────

def analyze_content(text: str) -> dict:
    try:
        raw = ai._generate(IMPORT_SYSTEM, f"Analyze this project content:\n\n{text[:10000]}", max_tokens=1500)
        match = re.search(r"\{.*\}", raw, re.DOTALL)
        return json.loads(match.group() if match else raw)
    except Exception:
        return _fallback_analysis()


def _fallback_analysis() -> dict:
    return {
        "project_name": "Imported Project",
        "description": "",
        "tags": [],
        "milestones": [],
        "import_narrative": "Content was imported. AI analysis unavailable — please edit project details manually.",
        "detected_tools": [],
        "health_score": 60,
        "health_explanation": "Score reflects that documentation was provided but analysis could not complete.",
    }


def url_to_link_dict(url: str) -> dict:
    from services import links_service
    detected = links_service.detect_platform(url)
    return {
        "platform": detected.get("platform", "other"),
        "label": detected.get("label", "Link"),
        "url": url,
        "icon": detected.get("icon", ""),
    }


# ── Project creation ───────────────────────────────────────────────────────────

def create_project_from_import(
    user_id: str,
    analysis: dict,
    raw_text: str,
    linked_urls: list[dict] | None = None,
) -> dict:
    project = db.create_project(
        owner_id=user_id,
        title=analysis.get("project_name", "Imported Project"),
        description=analysis.get("description", ""),
        tags=analysis.get("tags", []),
    )
    project_id = project["id"]

    supabase = get_supabase()
    supabase.table("projects").update({
        "health_score": analysis.get("health_score", 60),
        "health_explanation": analysis.get("health_explanation", ""),
    }).eq("id", project_id).execute()

    # Initial snapshot
    db.insert_snapshot({
        "project_id": project_id,
        "created_by": user_id,
        "title": "Initial Import Snapshot",
        "description": "Auto-generated on project import.",
        "ai_narrative": analysis.get("import_narrative", ""),
        "snapshot_data": {
            "project": db.get_project(project_id),
            "milestones": analysis.get("milestones", []),
            "source": "import",
        },
        "trigger": "auto",
    })

    # Milestones as updates
    for m in analysis.get("milestones", [])[:10]:
        db.add_update(
            project_id=project_id,
            user_id=user_id,
            content=f"Milestone: {m.get('title', '')} ({m.get('date', 'unknown date')})",
            update_type="milestone",
            ai_summary=m.get("title", ""),
        )

    # Imported content as a searchable note + embed
    _narrative = re.sub(r'<[^>]+>', '', str(analysis.get("import_narrative", "") or "")).strip()
    _narrative = ' '.join(_narrative.split())
    if raw_text.strip():
        _safe_content = re.sub(r'<[^>]+>', '', raw_text).strip()[:500]
        _safe_content = ' '.join(_safe_content.split())
        update = db.add_update(
            project_id=project_id,
            user_id=user_id,
            content=_safe_content,
            update_type="note",
            ai_summary=_narrative,
        )
        try:
            search_service.embed_update(update, project_id)
        except Exception:
            pass

    # Auto-add detected tools as linked workspaces
    from services import links_service
    for tool in analysis.get("detected_tools", []):
        if tool in links_service.PLATFORM_PRESETS:
            preset = links_service.PLATFORM_PRESETS[tool]
            if preset.get("domain"):
                links_service.add_project_link(
                    project_id=project_id,
                    user_id=user_id,
                    platform=tool,
                    label=preset["label"],
                    url=f"https://{preset['domain']}",
                    icon=preset["icon"],
                )

    # Add explicitly provided linked URLs (e.g. Google Drive files)
    for link in (linked_urls or []):
        if link.get("url"):
            links_service.add_project_link(
                project_id=project_id,
                user_id=user_id,
                platform=link.get("platform", "other"),
                label=link.get("label", "Imported File"),
                url=link["url"],
                icon=link.get("icon", ""),
            )

    return db.get_project(project_id)
